#!/usr/bin/env python3
"""
Generate pbmc_llm_pipeline_comparison.pptx
8 slides comparing Claude / Llama / Gemma pipelines.
Charts → embedded matplotlib PNG; tables & text → native pptx (editable).
"""

from io import BytesIO
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ── Constants ─────────────────────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)   # 16:9 widescreen

# Colours
C = dict(
    dark   = RGBColor(0x1A, 0x25, 0x2F),
    mid    = RGBColor(0x5D, 0x6D, 0x7E),
    light  = RGBColor(0xF8, 0xF9, 0xFA),
    claude = RGBColor(0x2E, 0x86, 0xC1),
    llama  = RGBColor(0xE6, 0x7E, 0x22),
    gemma  = RGBColor(0x27, 0xAE, 0x60),
    zs     = RGBColor(0x8E, 0x44, 0xAD),
    ref    = RGBColor(0xC0, 0x39, 0x2B),
    white  = RGBColor(0xFF, 0xFF, 0xFF),
    header = RGBColor(0x1B, 0x4F, 0x72),
)
HEX = {k: f"{v[0]:02X}{v[1]:02X}{v[2]:02X}" for k, v in C.items()}

MODEL_COLORS = ["#2E86C1", "#E67E22", "#27AE60"]   # matplotlib-compatible hex strings

# Metrics data
METRICS = {
    "Zero-Shot": {
        "Claude (Haiku 4.5)": dict(acc=0.851, mf1=0.673, wf1=0.804, kap=0.794, time=2.0,  mem=10),
        "Llama (3.1 70B)":    dict(acc=0.851, mf1=0.673, wf1=0.804, kap=0.794, time=1.1,  mem=80),
        "Gemma (3 27B)":      dict(acc=0.764, mf1=0.482, wf1=0.692, kap=0.674, time=23.5, mem=216),
    },
    "Reference Mapping": {
        "Claude (CellTypist)":   dict(acc=0.889, mf1=0.776, wf1=0.863, kap=0.847, time=1.5,  mem=255),
        "Llama (dot-product)":   dict(acc=0.732, mf1=0.502, wf1=0.710, kap=0.650, time=0.3,  mem=35),
        "Gemma (Scanorama+KNN)": dict(acc=0.770, mf1=0.564, wf1=0.760, kap=0.697, time=1.4,  mem=77),
    },
}
PER_CLASS = {
    "Zero-Shot":  {"Claude": pd.read_csv("results/evaluation/per_class_zero_shot.csv",     index_col=0),
                   "Llama":  pd.read_csv("results_llama/evaluation/per_class_zero_shot.csv", index_col=0),
                   "Gemma":  pd.read_csv("results_gemma/evaluation/per_class_zero_shot.csv", index_col=0)},
    "Reference":  {"Claude": pd.read_csv("results/evaluation/per_class_reference.csv",     index_col=0),
                   "Llama":  pd.read_csv("results_llama/evaluation/per_class_reference.csv", index_col=0),
                   "Gemma":  pd.read_csv("results_gemma/evaluation/per_class_reference.csv", index_col=0)},
}
CELL_TYPES = ["B", "CD4 T", "CD8 T", "CD14 Mono", "DC", "FCGR3A Mono", "NK", "Platelet"]

# ── pptx helpers ──────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, rgb: RGBColor):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = rgb

def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_w=0):
    shp = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb

def add_img(slide, fig, l, t, w, h, dpi=150):
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, l, t, w, h)
    plt.close(fig)

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", hex_color)

def cell_text(cell, text, size=10, bold=False, color="000000", align=PP_ALIGN.CENTER):
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    p.clear()
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.85), C["dark"])
    add_text(slide, title, Inches(0.25), Inches(0.05), Inches(12), Inches(0.5),
             size=24, bold=True, color=C["white"], align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.25), Inches(0.52), Inches(12.8), Inches(0.28),
                 size=11, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.LEFT)

def footer_bar(slide, n, total=8):
    add_rect(slide, 0, SH - Inches(0.28), SW, Inches(0.28), C["dark"])
    add_text(slide,
             f"PBMC Cell-Type Annotation: LLM Pipeline Comparison  |  Scripps Research Hackathon 2026  |  {n}/{total}",
             Inches(0.1), SH - Inches(0.27), SW - Inches(0.2), Inches(0.26),
             size=8, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)

# ── Chart helpers (matplotlib → embedded PNG) ─────────────────────────────────
def bar_chart(method_key, w_fig=6.5, h_fig=4.2):
    fig, ax = plt.subplots(figsize=(w_fig, h_fig))
    fig.patch.set_facecolor("#F8F9FA"); ax.set_facecolor("#F8F9FA")
    metric_keys  = ["acc", "mf1", "wf1", "kap"]
    metric_names = ["Accuracy", "Macro F1", "Weighted F1", "Cohen's κ"]
    models = list(METRICS[method_key].keys())
    x = np.arange(len(metric_names)); bw = 0.25
    for i, (model, color) in enumerate(zip(models, MODEL_COLORS)):
        vals = [METRICS[method_key][model][k] for k in metric_keys]
        bars = ax.bar(x + (i-1)*bw, vals, bw,
                      label=model, color=color, alpha=0.88,
                      edgecolor="white", linewidth=0.5)
        for bar, val in zip(bars, vals):
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.01,
                    f"{val:.2f}", ha="center", va="bottom", fontsize=7.5, fontweight="bold")
    ax.set_xticks(x); ax.set_xticklabels(metric_names, fontsize=11)
    ax.set_ylim(0, 1.18); ax.set_ylabel("Score", fontsize=11)
    ax.legend(fontsize=9); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout()
    return fig

def perclass_chart(w_fig=12, h_fig=4.5):
    fig, axes = plt.subplots(1, 2, figsize=(w_fig, h_fig))
    fig.patch.set_facecolor("#F8F9FA")
    for ax, (mk, mlabel) in zip(axes, [("Zero-Shot","Zero-Shot"),("Reference","Reference Mapping")]):
        ax.set_facecolor("#F8F9FA")
        x = np.arange(len(CELL_TYPES)); bw = 0.26
        for i, (name, color) in enumerate(zip(["Claude","Llama","Gemma"], MODEL_COLORS)):
            vals = [float(PER_CLASS[mk][name].loc[ct, "f1-score"])
                    if ct in PER_CLASS[mk][name].index else 0.0 for ct in CELL_TYPES]
            ax.bar(x+(i-1)*bw, vals, bw, label=name, color=color, alpha=0.85, edgecolor="white")
        ax.set_xticks(x); ax.set_xticklabels(CELL_TYPES, rotation=35, ha="right", fontsize=9)
        ax.set_ylim(0, 1.15); ax.set_ylabel("F1 Score", fontsize=10)
        ax.set_title(mlabel, fontsize=12, fontweight="bold")
        ax.legend(fontsize=9); ax.spines[["top","right"]].set_visible(False)
        ax.axhspan(0, 0.05, alpha=0.08, color="red")
    fig.tight_layout()
    return fig

def resource_chart(w_fig=11, h_fig=4.0):
    labels = ["ZS-Claude","ZS-Llama","ZS-Gemma","Ref-Claude","Ref-Llama","Ref-Gemma"]
    times  = [2.0, 1.1, 23.5, 1.5, 0.3, 1.4]
    mems   = [10, 80, 216, 255, 35, 77]
    colors = [c for c in MODEL_COLORS*2]
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(w_fig, h_fig))
    fig.patch.set_facecolor("#F8F9FA")
    for ax, vals, xlabel, title in [
        (ax1, times, "Wall Time (s)", "Runtime per Method"),
        (ax2, mems,  "Peak Memory (MB)", "Peak Memory per Method"),
    ]:
        ax.set_facecolor("#F8F9FA")
        bars = ax.barh(labels, vals, color=colors, alpha=0.85, edgecolor="white")
        for bar, val in zip(bars, vals):
            ax.text(val + 0.5, bar.get_y()+bar.get_height()/2,
                    f"{val}" + ("s" if xlabel.startswith("Wall") else " MB"),
                    va="center", fontsize=9, fontweight="bold")
        ax.set_xlabel(xlabel, fontsize=10); ax.set_title(title, fontsize=11, fontweight="bold")
        ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout()
    return fig

# ══════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── Slide 1: Title ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["dark"])
add_rect(sl, 0, Inches(2.3), SW, Inches(3.2), C["header"])

add_text(sl, "PBMC Cell-Type Annotation",
         Inches(0.5), Inches(2.55), Inches(12.3), Inches(1.0),
         size=40, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
add_text(sl, "LLM Pipeline Comparison: Claude vs Llama vs Gemma",
         Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.7),
         size=22, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)
add_text(sl, "Dataset: PBMC 3k  |  AWS Bedrock  |  scanpy · CellTypist · Scanorama",
         Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.4),
         size=13, color=C["mid"], align=PP_ALIGN.CENTER)

for i, (label, color) in enumerate([("Claude Haiku 4.5", C["claude"]),
                                      ("Llama 3.1 70B",    C["llama"]),
                                      ("Gemma 3 27B",      C["gemma"])]):
    x = Inches(3.2 + i * 2.5)
    add_rect(sl, x, Inches(5.6), Inches(2.1), Inches(0.65), color)
    add_text(sl, label, x, Inches(5.62), Inches(2.1), Inches(0.62),
             size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

add_text(sl, "Scripps Research Hackathon 2026",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         size=11, color=C["mid"], align=PP_ALIGN.CENTER)
footer_bar(sl, 1)

# ── Slide 2: Pipeline Design Comparison Table ──────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Pipeline Design Comparison",
           "What did each LLM choose for each method?")

rows_data = [
    ["", "Claude (Haiku 4.5)", "Llama (3.1 70B)", "Gemma (3 27B)"],
    ["Zero-Shot Approach",    "LLM API call via Bedrock\n(biological reasoning)",
                              "Avg marker gene expression\nper cluster",
                              "Sum marker gene expression\nper cell"],
    ["Zero-Shot Granularity", "Cluster-level",        "Cluster-level",       "Per-cell"],
    ["Reference Method",      "CellTypist\n(logistic regression)",
                              "Centroid dot-product\nsimilarity",
                              "Scanorama integration\n+ KNN label transfer"],
    ["Reference Dataset",     "Immune_All_High\n(700k cells, 20+ studies)",
                              "pbmc68k_reduced\n(700 cells, 10x)",
                              "pbmc68k_reduced\n(700 cells, 10x)"],
    ["Code Structure",        "4 scripts (pipeline)",  "1 monolithic script", "6 modular scripts"],
    ["Requires Bedrock API",  "Yes (LLM call)",        "No",                  "No"],
    ["Accuracy — Zero-Shot",  "85.1%", "85.1%", "76.4%"],
    ["Accuracy — Reference",  "88.9%", "73.2%", "77.0%"],
]

tbl_left = Inches(0.25); tbl_top = Inches(0.95)
tbl_w = Inches(12.83);   tbl_h = Inches(6.3)
tbl_shape = sl.shapes.add_table(len(rows_data), 4, tbl_left, tbl_top, tbl_w, tbl_h)
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.3)
for c in range(1, 4): tbl.columns[c].width = Inches(3.51)

ROW_BG = ["1B4F72","D6EAF8","E8F8F5","D6EAF8","E8F8F5","D6EAF8","E8F8F5","D6EAF8","E8F8F5"]
for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        if r == 0:
            bg = "1B4F72" if c == 0 else ["1B4F72","2E86C1","E67E22","27AE60"][c]
            set_cell_bg(cell, bg)
            cell_text(cell, val, size=11, bold=True, color="FFFFFF")
        else:
            col_bg = ["D5D8DC","D6EAF8","FDEBD0","D5F5E3"][c]
            bg = "EAECEE" if r%2==0 and c==0 else ("D5D8DC" if c==0 else col_bg)
            if r%2==0 and c>0:
                bg = ["F2F3F4","EBF5FB","FEF9E7","EAFAF1"][c]
            set_cell_bg(cell, bg)
            cell_text(cell, val, size=10, color="1A252F")

footer_bar(sl, 2)

# ── Slide 3: Zero-Shot Results ─────────────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Zero-Shot Annotation Results",
           "Cluster marker genes → cell type label  |  No reference data used")

fig = bar_chart("Zero-Shot")
add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(7.2), Inches(5.0))

# Summary table (right)
sum_rows = [["Model", "Accuracy", "Macro F1", "Wtd F1", "Kappa", "Time", "Mem"],
            ["Claude (Haiku 4.5)", "0.851", "0.673", "0.804", "0.794", "2.0 s", "10 MB"],
            ["Llama (3.1 70B)",    "0.851", "0.673", "0.804", "0.794", "1.1 s", "80 MB"],
            ["Gemma (3 27B)",      "0.764", "0.482", "0.692", "0.674", "23.5 s","216 MB"]]
ts = sl.shapes.add_table(4, 7, Inches(7.7), Inches(1.5), Inches(5.5), Inches(2.0))
tt = ts.table
for c, w in enumerate([1.5,0.75,0.75,0.75,0.75,0.75,0.75]):
    tt.columns[c].width = Inches(w)
for r, row in enumerate(sum_rows):
    for c, val in enumerate(row):
        cell = tt.cell(r, c)
        if r == 0:
            set_cell_bg(cell, "1B4F72")
            cell_text(cell, val, size=9, bold=True, color="FFFFFF")
        else:
            bg = ["D6EAF8","FDEBD0","D5F5E3"][r-1]
            set_cell_bg(cell, bg)
            cell_text(cell, val, size=9)

add_text(sl, "Claude & Llama tie at 85.1% — identical cluster assignments\ndespite completely different mechanisms",
         Inches(7.7), Inches(3.7), Inches(5.5), Inches(0.8),
         size=11, color=C["claude"], bold=True, italic=True, align=PP_ALIGN.CENTER)

footer_bar(sl, 3)

# ── Slide 4: Reference Mapping Results ────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Reference Mapping Results",
           "Annotate by comparing against a large, pre-labelled cell atlas")

fig = bar_chart("Reference Mapping")
add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(7.2), Inches(5.0))

sum_rows = [["Model", "Accuracy", "Macro F1", "Wtd F1", "Kappa", "Time", "Mem"],
            ["Claude (CellTypist)",   "0.889", "0.776", "0.863", "0.847", "1.5 s", "255 MB"],
            ["Llama (dot-product)",   "0.732", "0.502", "0.710", "0.650", "0.3 s", "35 MB"],
            ["Gemma (Scanorama+KNN)", "0.770", "0.564", "0.760", "0.697", "1.4 s", "77 MB"]]
ts = sl.shapes.add_table(4, 7, Inches(7.7), Inches(1.5), Inches(5.5), Inches(2.0))
tt = ts.table
for c, w in enumerate([1.6,0.7,0.7,0.7,0.7,0.7,0.7]):
    tt.columns[c].width = Inches(w)
for r, row in enumerate(sum_rows):
    for c, val in enumerate(row):
        cell = tt.cell(r, c)
        if r == 0:
            set_cell_bg(cell, "1B4F72"); cell_text(cell, val, size=9, bold=True, color="FFFFFF")
        else:
            bg = ["D6EAF8","FDEBD0","D5F5E3"][r-1]
            set_cell_bg(cell, bg); cell_text(cell, val, size=9)

add_text(sl, "CellTypist (Claude) wins by +12-16 pp\nby choosing a pre-trained 700k-cell model",
         Inches(7.7), Inches(3.7), Inches(5.5), Inches(0.8),
         size=11, color=C["claude"], bold=True, italic=True, align=PP_ALIGN.CENTER)

footer_bar(sl, 4)

# ── Slide 5: Per-Class F1 ─────────────────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Per-Class F1 Score Breakdown",
           "Which cell types does each method handle well — and which does it miss?")

fig = perclass_chart()
add_img(sl, fig, Inches(0.3), Inches(0.95), Inches(12.7), Inches(5.7))

add_text(sl, "Red zone (F1≈0): FCGR3A Mono undetected by all — merged with CD14 Mono in the same Leiden cluster",
         Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.35),
         size=10, color=C["ref"], bold=True, italic=True, align=PP_ALIGN.CENTER)
footer_bar(sl, 5)

# ── Slide 6: UMAP ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "UMAP Visualisations — Claude Pipeline",
           "Ground Truth (left)  |  Zero-Shot LLM (centre)  |  CellTypist Reference Mapping (right)")

umap_path = "results/evaluation/umap_three_way.png"
if Path(umap_path).exists():
    sl.shapes.add_picture(umap_path, Inches(0.2), Inches(0.95), Inches(12.9), Inches(5.8))
else:
    add_text(sl, "UMAP image not found — run 04_evaluate.py first",
             Inches(1), Inches(3), Inches(11), Inches(1), size=14, color=C["mid"],
             align=PP_ALIGN.CENTER)
footer_bar(sl, 6)

# ── Slide 7: Resource Usage ───────────────────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Resource Usage: Time & Memory",
           "Wall-clock time and peak memory for each annotation step")

fig = resource_chart()
add_img(sl, fig, Inches(0.4), Inches(0.95), Inches(12.5), Inches(5.5))
add_text(sl, "Lighter = Zero-Shot   |   Solid = Reference Mapping",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.35),
         size=10, color=C["mid"], italic=True, align=PP_ALIGN.CENTER)
footer_bar(sl, 7)

# ── Slide 8: Key Differences & Limitations ────────────────────────────────────
sl = blank_slide(prs)
set_bg(sl, C["light"])
header_bar(sl, "Key Differences & Limitations")

sections = [
    (C["claude"], "[1] METHOD DIFFERENCES", [
        "Claude: Only LLM that called the Bedrock API — used biological reasoning to annotate from marker genes",
        "Llama: Cluster-level marker scoring — no API call, faster, but identical zero-shot accuracy to Claude",
        "Gemma: Suggested Scanorama (proper batch-correction) — most sophisticated integration approach",
        "Reference tool choice drove the biggest accuracy gap: CellTypist (700k cells) vs dot-product/KNN (700 cells)",
    ]),
    (C["ref"], "[2] SHARED LIMITATIONS", [
        "Leiden resolution 0.5 yields 6 clusters for 8 ground-truth types — FCGR3A Mono & NK merge with neighbors",
        "FCGR3A+ Monocytes: F1 = 0.00 for ALL methods — requires higher clustering resolution to rescue",
        "pbmc68k reference (Llama/Gemma) too small (700 cells) for reliable centroid estimation",
        "Ground truth labels from scanpy are coarser than Seurat (no Naive/Memory CD4 T distinction)",
    ]),
    (C["llama"], "[3] LLM CODE QUALITY", [
        "Claude: Produced a complete, runnable pipeline — chose CellTypist without being prompted to",
        "Llama: Skeleton with bugs (wrong argmax logic, wrong boto3 syntax) — required ~10 fixes to run",
        "Gemma: Most structured output (6 files) but wrong reference (suggested Mouse Brain Atlas for PBMC!)",
        "All 3 omitted memory/time tracking — left as 'exercise for the reader'",
    ]),
]

y = Inches(1.0)
for color, title, bullets in sections:
    add_rect(sl, Inches(0.25), y, Inches(12.83), Inches(0.42), color)
    add_text(sl, title, Inches(0.35), y + Inches(0.04), Inches(12.5), Inches(0.38),
             size=12, bold=True, color=C["white"], align=PP_ALIGN.LEFT)
    y += Inches(0.44)
    for bullet in bullets:
        add_text(sl, f"•   {bullet}",
                 Inches(0.4), y, Inches(12.5), Inches(0.38),
                 size=10, color=RGBColor(0x1A, 0x25, 0x2F), align=PP_ALIGN.LEFT)
        y += Inches(0.39)
    y += Inches(0.05)

footer_bar(sl, 8)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "pbmc_llm_pipeline_comparison.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({Path(OUT).stat().st_size//1024} KB)")
