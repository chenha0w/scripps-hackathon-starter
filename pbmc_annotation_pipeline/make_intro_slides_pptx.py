#!/usr/bin/env python3
"""
Generate intro_cell_type_annotation.pptx
10 educational slides for a mixed biology/chemistry audience.
Diagrams → embedded matplotlib PNG; tables & text → native pptx (editable).
"""

from io import BytesIO
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import scipy.sparse as scipy_sparse
import scanpy as sc
import warnings
warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Load real PBMC data ───────────────────────────────────────────────────────
adata = sc.read_h5ad("results/pbmc_processed.h5ad")
adata_raw = adata.raw.to_adata() if adata.raw else adata

# ── Design tokens ─────────────────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)

C = dict(
    dark    = RGBColor(0x1A, 0x25, 0x2F),
    mid     = RGBColor(0x5D, 0x6D, 0x7E),
    light   = RGBColor(0xF8, 0xF9, 0xFA),
    white   = RGBColor(0xFF, 0xFF, 0xFF),
    blue    = RGBColor(0x2E, 0x86, 0xC1),
    green   = RGBColor(0x27, 0xAE, 0x60),
    purple  = RGBColor(0x6C, 0x34, 0x83),
    orange  = RGBColor(0xE6, 0x7E, 0x22),
    red     = RGBColor(0xE7, 0x4C, 0x3C),
    navy    = RGBColor(0x1B, 0x4F, 0x72),
    teal    = RGBColor(0x1A, 0x52, 0x76),
    lblue   = RGBColor(0xAE, 0xD6, 0xF1),
)
CT_PAL = {
    "CD4 T cells":        "#3498DB", "CD8 T cells":        "#1ABC9C",
    "NK cells":           "#9B59B6", "B cells":            "#E74C3C",
    "CD14+ Monocytes":   "#E67E22", "FCGR3A+ Monocytes":  "#F39C12",
    "Dendritic cells":   "#2ECC71", "Megakaryocytes":      "#95A5A6",
}
LEIDEN_PAL = {"0":"#3498DB","1":"#E67E22","2":"#9B59B6",
               "3":"#E74C3C","4":"#2ECC71","5":"#95A5A6"}

# ── pptx helpers ──────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    return prs

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, rgb):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else setattr(s.line.color, "rgb", line)
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color
    return tb

def img(slide, fig, l, t, w, h, dpi=150):
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0); slide.shapes.add_picture(buf, l, t, w, h); plt.close(fig)

def set_cell_bg(cell, hexc):
    tc = cell._tc; pr = tc.get_or_add_tcPr()
    for o in pr.findall(qn("a:solidFill")): pr.remove(o)
    sf = etree.SubElement(pr, qn("a:solidFill"))
    sc_el = etree.SubElement(sf, qn("a:srgbClr")); sc_el.set("val", hexc)

def ctxt(cell, text, size=10, bold=False, color="1A252F", align=PP_ALIGN.CENTER):
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; p.clear()
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)

def header(slide, title, subtitle=None, accent=None):
    accent = accent or C["dark"]
    rect(slide, 0, 0, SW, Inches(0.85), accent)
    txt(slide, title, Inches(0.25), Inches(0.05), Inches(12.5), Inches(0.5),
        size=22, bold=True, color=C["white"], align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, Inches(0.25), Inches(0.52), Inches(12.8), Inches(0.3),
            size=10.5, color=C["lblue"], align=PP_ALIGN.LEFT)

def footer(slide, n, total=10):
    rect(slide, 0, SH-Inches(0.28), SW, Inches(0.28), C["dark"])
    txt(slide,
        f"Introduction to Cell-Type Annotation  |  Scripps Research Hackathon 2026  |  {n}/{total}",
        Inches(0.1), SH-Inches(0.27), SW-Inches(0.2), Inches(0.26),
        size=8, color=C["lblue"], align=PP_ALIGN.CENTER)

def callout(slide, text, l, t, w, h, fill_rgb, text_rgb=None):
    rect(slide, l, t, w, h, fill_rgb)
    txt(slide, text, l+Inches(0.08), t+Inches(0.06), w-Inches(0.16), h-Inches(0.08),
        size=10.5, color=text_rgb or C["white"], align=PP_ALIGN.CENTER, bold=True)

# ── matplotlib figure generators ──────────────────────────────────────────────
def fig_central_question():
    fig, ax = plt.subplots(figsize=(12, 5.5))
    fig.patch.set_facecolor("#FAFAFA"); ax.set_facecolor("#FAFAFA")
    ax.set_xlim(0,10); ax.set_ylim(0,5); ax.axis("off")

    tissue = mpatches.Ellipse((2,2.5),3,4,facecolor="#FADBD8",edgecolor="#E74C3C",lw=2,alpha=0.5)
    ax.add_patch(tissue)
    ax.text(2,4.75,"Blood sample",ha="center",fontsize=12,fontweight="bold",color="#1A252F")

    np.random.seed(42)
    cols = ["#3498DB","#E74C3C","#E67E22","#9B59B6","#2ECC71"]
    for _ in range(45):
        cx,cy = np.random.uniform(0.6,3.4),np.random.uniform(0.6,4.4)
        c = mpatches.Circle((cx,cy),0.22,facecolor=cols[np.random.randint(5)],edgecolor="white",lw=1,alpha=0.85)
        ax.add_patch(c)
    ax.text(2,0.1,"All cells look the same\nto the naked eye",ha="center",fontsize=10,color="#5D6D7E",style="italic")

    ax.annotate("",xy=(5.3,2.5),xytext=(3.7,2.5),
                arrowprops=dict(arrowstyle="-|>",color="#1A252F",lw=2.5,mutation_scale=22))
    ax.text(4.5,3.0,"scRNA-seq\n+\nannotation",ha="center",fontsize=11,fontweight="bold",color="#2E86C1")

    labels=[("T cells","#3498DB"),("B cells","#E74C3C"),("Monocytes","#E67E22"),("NK cells","#9B59B6"),("DC","#2ECC71")]
    for i,(lbl,col) in enumerate(labels):
        y=4.2-i*0.85
        for j in range(6):
            c=mpatches.Circle((6.1+j*0.42,y),0.18,facecolor=col,edgecolor="white",lw=1,alpha=0.9)
            ax.add_patch(c)
        ax.text(9.0,y,lbl,va="center",fontsize=11,fontweight="bold",color=col)
    ax.text(7.5,4.85,"Cells sorted by identity!",ha="center",fontsize=12,fontweight="bold",color="#1A252F")
    ax.text(7.5,0.1,"Annotation reveals what each cell actually IS",ha="center",
            fontsize=11,color="#2E86C1",fontweight="bold")
    fig.tight_layout(); return fig

def fig_central_dogma():
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.patch.set_facecolor("#FAFAFA"); ax.set_facecolor("#FAFAFA")
    ax.set_xlim(0,10); ax.set_ylim(0,5); ax.axis("off")

    steps = [(1.2,3.5,"#2C3E50","DNA","Blueprint\nstored in\nevery cell"),
             (4.0,3.5,"#8E44AD","mRNA","Active copy\nread from\nthe DNA"),
             (6.8,3.5,"#E67E22","Protein","Functional\nmolecule that\ndoes the work")]
    for x,y,col,t,d in steps:
        circ=mpatches.FancyBboxPatch((x-0.9,y-0.6),1.8,1.2,boxstyle="round,pad=0.08",
                                      facecolor=col,edgecolor="white",lw=2)
        ax.add_patch(circ)
        ax.text(x,y+0.2,t,ha="center",va="center",fontsize=14,fontweight="bold",color="white")
        ax.text(x,y-1.1,d,ha="center",va="top",fontsize=9.5,color="#5D6D7E",multialignment="center")
    for x0,x1,lbl in [(2.15,2.95,"transcription"),(4.95,5.75,"translation")]:
        ax.annotate("",xy=(x1,3.5),xytext=(x0,3.5),
                    arrowprops=dict(arrowstyle="-|>",color="#1A252F",lw=2,mutation_scale=16))
        ax.text((x0+x1)/2,3.78,lbl,ha="center",fontsize=9,color="#5D6D7E")
    ax.text(5,4.75,"The Central Dogma of Molecular Biology",ha="center",
            fontsize=12,fontweight="bold",color="#1A252F")

    ax.axhline(2.25,color="#D5D8DC",lw=1.5,ls="--")
    ax.text(0.15,2.0,"Chemistry analogy:",fontsize=11,fontweight="bold",color="#6C3483")
    ax.text(0.15,0.3,"scRNA-seq is like a chromatogram for each cell:\ninstead of measuring molecular masses, we count how often\neach gene's mRNA is detected. High count = gene is 'active'.",
            fontsize=10.5,color="#1A252F",va="top")

    box=mpatches.FancyBboxPatch((7.0,0.15),2.85,2.05,boxstyle="round,pad=0.08",
                                 facecolor="#EBF5FB",edgecolor="#2E86C1",lw=2)
    ax.add_patch(box)
    ax.text(8.42,1.95,"Key insight",ha="center",fontsize=11,fontweight="bold",color="#2E86C1")
    ax.text(8.42,1.15,"Different cell types use\ndifferent sets of genes.\nThis IS their identity.",
            ha="center",va="center",fontsize=10.5,color="#1A252F",multialignment="center")
    fig.tight_layout(); return fig

def fig_workflow():
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.patch.set_facecolor("#FAFAFA"); ax.set_facecolor("#FAFAFA")
    ax.set_xlim(0,10); ax.set_ylim(0,5); ax.axis("off")
    steps=[
        (0.9,2.6,"#E8DAEF","#8E44AD","Tissue\nSample","Blood / tumor\norgan biopsy"),
        (2.7,2.6,"#D6EAF8","#2E86C1","Cell\nDissociation","Break into\nsingle cells"),
        (4.5,2.6,"#D5F5E3","#27AE60","Droplet\nCapture","10x Chromium\nbarcodes cells"),
        (6.3,2.6,"#FDEBD0","#E67E22","Sequencing","Read mRNA\nper droplet"),
        (8.1,2.6,"#FDEDEC","#E74C3C","Count\nMatrix","Cells x Genes\ntable"),
    ]
    for i,(x,y,fc,ec,t,d) in enumerate(steps):
        c=mpatches.Circle((x,y),0.75,facecolor=fc,edgecolor=ec,lw=2.5)
        ax.add_patch(c)
        ax.text(x,y+0.18,t,ha="center",va="center",fontsize=9.5,fontweight="bold",color=ec,multialignment="center")
        ax.text(x,y-1.2,f"Step {i+1}",ha="center",fontsize=9,fontweight="bold",color=ec)
        ax.text(x,y-1.55,d,ha="center",va="top",fontsize=8.5,color="#5D6D7E",multialignment="center")
        if i<4:
            ax.annotate("",xy=(x+0.95,y),xytext=(x+0.82,y),
                        arrowprops=dict(arrowstyle="-|>",color="#1A252F",lw=2,mutation_scale=16))

    # droplet inset
    da=fig.add_axes([0.36,0.65,0.12,0.28])
    da.set_xlim(0,3); da.set_ylim(0,3); da.axis("off")
    da.add_patch(mpatches.Circle((1.5,1.5),1.3,facecolor="#D5F5E3",edgecolor="#27AE60",lw=2,alpha=0.7))
    da.add_patch(mpatches.Circle((1.5,1.5),0.55,facecolor="#AED6F1",edgecolor="#2E86C1",lw=1.5))
    da.text(1.5,1.5,"cell",ha="center",va="center",fontsize=8,fontweight="bold",color="#2E86C1")
    for dy in [0.85,1.05,1.25]:
        for dx in [-0.6,0.6]:
            da.plot([1.5+dx*0.6,1.5+dx*0.9],[1.5+dy*0.5,1.5+dy*0.5],color="#E74C3C",lw=1.5,alpha=0.7)
    da.text(1.5,0.1,"+ barcode",ha="center",fontsize=7,color="#5D6D7E")

    # mini heatmap
    np.random.seed(7)
    mat=np.random.randint(0,30,(5,5)).astype(float)
    mat[0,:2]=[48,40]; mat[1,2:4]=[44,36]; mat[2,1]=50
    ma=fig.add_axes([0.812,0.12,0.16,0.36])
    ma.imshow(mat,cmap="YlOrRd",aspect="auto")
    ma.set_xticks(range(5)); ma.set_xticklabels(["CD3D","LYZ","CD79A","NKG7","PPBP"],rotation=45,ha="right",fontsize=6.5)
    ma.set_yticks(range(5)); ma.set_yticklabels([f"Cell {i+1}" for i in range(5)],fontsize=6.5)
    ma.set_title("Count matrix",fontsize=8,fontweight="bold",color="#E74C3C")
    fig.tight_layout(); return fig

def fig_umap_panels():
    fig, axes = plt.subplots(1,3,figsize=(13,4.5))
    fig.patch.set_facecolor("#FAFAFA")
    umap = adata.obsm["X_umap"]
    leiden = adata.obs["leiden"].values
    gt = adata.obs["ground_truth"].values

    # PCA variance
    ax=axes[0]; ax.set_facecolor("#FAFAFA")
    var=[7.8,6.1,4.4,3.8,2.9,2.4,2.0,1.8,1.6,1.4]
    ax.bar(range(1,11),var,color="#2E86C1",edgecolor="white",alpha=0.85)
    ax.set_xlabel("Principal Component",fontsize=10); ax.set_ylabel("Variance (%)",fontsize=10)
    ax.set_title("Step 1: PCA\n(compress 20k genes → 50 PCs)",fontsize=10,fontweight="bold")
    ax.spines[["top","right"]].set_visible(False)

    # Unlabelled UMAP
    ax=axes[1]; ax.set_facecolor("#FAFAFA")
    for cl in sorted(set(leiden)):
        m=leiden==cl
        ax.scatter(umap[m,0],umap[m,1],c=LEIDEN_PAL[cl],s=3,alpha=0.65,rasterized=True)
        cx,cy=umap[m,0].mean(),umap[m,1].mean()
        ax.text(cx,cy,cl,ha="center",va="center",fontsize=13,fontweight="bold",color="white",
                path_effects=[pe.withStroke(linewidth=3,foreground="black")])
    ax.set_title("Step 2: UMAP\n(unlabelled clusters)",fontsize=10,fontweight="bold")
    ax.set_xlabel("UMAP 1"); ax.set_ylabel("UMAP 2"); ax.spines[["top","right"]].set_visible(False)
    ax.text(0.5,-0.14,"Each dot = one cell",ha="center",transform=ax.transAxes,fontsize=9,color="#5D6D7E")

    # Labelled UMAP
    ax=axes[2]; ax.set_facecolor("#FAFAFA")
    for ct,col in CT_PAL.items():
        m=gt==ct
        ax.scatter(umap[m,0],umap[m,1],c=col,s=3,alpha=0.7,label=ct,rasterized=True)
    ax.set_title("Step 3: Annotated!\n(cell types revealed)",fontsize=10,fontweight="bold",color="#27AE60")
    ax.set_xlabel("UMAP 1"); ax.set_ylabel("UMAP 2"); ax.spines[["top","right"]].set_visible(False)
    ax.legend(fontsize=6,markerscale=3,loc="lower right",framealpha=0.9)
    ax.text(0.5,-0.14,"Goal of annotation",ha="center",transform=ax.transAxes,
            fontsize=9,color="#27AE60",fontweight="bold")
    fig.tight_layout(); return fig

def fig_dotplot():
    markers_dict = {
        "B cells":["CD79A","MS4A1"],"CD4 T cells":["CD3D","IL7R"],
        "CD8 T cells":["CD8A","NKG7"],"NK cells":["GNLY","NKG7"],
        "CD14+ Monocytes":["CD14","LYZ"],"FCGR3A+ Monocytes":["FCGR3A","MS4A7"],
        "Dendritic cells":["FCER1A","CST3"],"Megakaryocytes":["PPBP","PF4"],
    }
    all_m=["CD79A","MS4A1","CD3D","IL7R","CD8A","GNLY","NKG7","CD14","LYZ","FCGR3A","MS4A7","FCER1A","CST3","PPBP","PF4"]
    cts=list(CT_PAL.keys())
    fig,ax=plt.subplots(figsize=(11,4.5))
    fig.patch.set_facecolor("#FAFAFA"); ax.set_facecolor("#FAFAFA")
    ct_i={ct:i for i,ct in enumerate(cts)}; g_i={g:i for i,g in enumerate(all_m)}
    for ct in cts:
        m=adata_raw.obs["ground_truth"]==ct
        if not m.any(): continue
        for gene in all_m:
            if gene not in adata_raw.var_names: continue
            X=adata_raw[m,gene].X
            arr=X.toarray().flatten() if scipy_sparse.issparse(X) else np.array(X).flatten()
            pct=100*(arr>0).mean(); mean=arr.mean()
            ax.scatter(g_i[gene],ct_i[ct],s=max(4,pct*2.5),
                       c=CT_PAL[ct],alpha=min(1.0,max(0.15,mean/3)),edgecolors="none")
    ax.set_xticks(range(len(all_m))); ax.set_xticklabels(all_m,rotation=45,ha="right",fontsize=8.5)
    ax.set_yticks(range(len(cts))); ax.set_yticklabels(cts,fontsize=9)
    for tick,ct in zip(ax.get_yticklabels(),cts):
        tick.set_color(CT_PAL[ct]); tick.set_fontweight("bold")
    ax.set_title("Dot Plot: size = % cells expressing  |  colour saturation = mean expression level",
                 fontsize=10,fontweight="bold")
    ax.grid(True,alpha=0.2); ax.spines[["top","right"]].set_visible(False)
    fig.tight_layout(); return fig

def fig_zero_shot_flow():
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.patch.set_facecolor("#FAFAFA"); ax.set_facecolor("#FAFAFA")
    ax.set_xlim(0,10); ax.set_ylim(0,5.5); ax.axis("off")

    # Top concept boxes
    for x,fc,ec,txt_s in [(2.2,"#F4ECF7","#8E44AD",'"Zero-shot" = no training data needed'),
                           (7.5,"#EBF5FB","#2E86C1","Works like a biologist looking up a reference book")]:
        b=mpatches.FancyBboxPatch((x-2.0,4.4),4.0,0.85,boxstyle="round,pad=0.06",facecolor=fc,edgecolor=ec,lw=2)
        ax.add_patch(b); ax.text(x,4.82,txt_s,ha="center",va="center",fontsize=10.5,fontweight="bold",color=ec)

    # Pipeline steps
    for i,(x,fc,ec,t) in enumerate([(1.0,"#D7BDE2","#8E44AD","UMAP\nclusters\n(unlabelled)"),
                                      (3.2,"#D6EAF8","#2E86C1","Find top\nmarker genes\nper cluster"),
                                      (5.4,"#FDEBD0","#E67E22","Match to\nknown marker\ndictionary / LLM"),
                                      (7.6,"#D5F5E3","#27AE60","Clusters\nlabelled with\ncell type")]):
        b=mpatches.FancyBboxPatch((x-0.85,2.35),1.7,1.3,boxstyle="round,pad=0.07",facecolor=fc,edgecolor=ec,lw=2)
        ax.add_patch(b); ax.text(x,3.0,t,ha="center",va="center",fontsize=9.5,fontweight="bold",color=ec,multialignment="center")
        if i<3:
            ax.annotate("",xy=(x+1.0,3.0),xytext=(x+0.9,3.0),
                        arrowprops=dict(arrowstyle="-|>",color="#1A252F",lw=2,mutation_scale=16))

    # Example
    ax.text(5,2.05,"Example: Cluster 3 top genes → CD79A, MS4A1, CD79B",
            ha="center",fontsize=10.5,fontweight="bold",color="#1A252F")
    for i,(x,fc,ec,t) in enumerate([(1.5,"#FDFEFE","#95A5A6","Cluster 3\ntop genes:\nCD79A, MS4A1"),
                                      (4.0,"#EBF5FB","#2E86C1","Known markers:\nCD79A = B cells\nMS4A1 = B cells"),
                                      (6.8,"#D5F5E3","#27AE60",'Annotation:\n"B cells"\nConfidence: HIGH')]):
        b=mpatches.FancyBboxPatch((x-1.0,0.2),2.0,1.5,boxstyle="round,pad=0.05",facecolor=fc,edgecolor=ec,lw=1.5)
        ax.add_patch(b); ax.text(x,0.95,t,ha="center",va="center",fontsize=9.5,color=ec,fontweight="bold",multialignment="center")
        if i<2: ax.annotate("",xy=(x+1.15,0.95),xytext=(x+1.05,0.95),
                             arrowprops=dict(arrowstyle="-|>",color="#1A252F",lw=1.5,mutation_scale=14))

    # LLM note
    b=mpatches.FancyBboxPatch((8.1,0.15),1.8,1.85,boxstyle="round,pad=0.05",facecolor="#FEF9E7",edgecolor="#F39C12",lw=2)
    ax.add_patch(b)
    ax.text(9.0,1.15,"LLM variant:\nSend marker genes\nto Claude/Llama/\nGemma via Bedrock",
            ha="center",va="center",fontsize=8.5,color="#E67E22",fontweight="bold",multialignment="center")
    fig.tight_layout(); return fig

def fig_reference_map():
    fig = plt.figure(figsize=(12, 5))
    fig.patch.set_facecolor("#FAFAFA")
    ax = fig.add_axes([0,0,1,1]); ax.set_xlim(0,10); ax.set_ylim(0,5.5); ax.axis("off")

    # GPS analogy banner
    b=mpatches.FancyBboxPatch((0.1,4.7),9.8,0.65,boxstyle="round,pad=0.06",facecolor="#EBF5FB",edgecolor="#2E86C1",lw=2)
    ax.add_patch(b)
    ax.text(5,5.03,"Analogy: Like GPS navigation — your new cells are an unknown location, the reference atlas is the map",
            ha="center",va="center",fontsize=10.5,fontweight="bold",color="#2E86C1")

    np.random.seed(1)
    ref_types=[("T cells","#3498DB",(1.1,2.8)),("B cells","#E74C3C",(2.4,1.5)),
               ("Monocytes","#E67E22",(0.7,1.1)),("NK cells","#9B59B6",(2.0,3.5)),("DC","#2ECC71",(2.9,2.2))]
    ref_x,ref_y,ref_c=[],[],[]
    for lbl,col,ctr in ref_types:
        xs=np.random.normal(ctr[0],0.22,60); ys=np.random.normal(ctr[1],0.22,60)
        ref_x.extend(xs); ref_y.extend(ys); ref_c.extend([col]*60)

    # Reference atlas panel
    ra=fig.add_axes([0.04,0.1,0.28,0.56]); ra.scatter(ref_x,ref_y,c=ref_c,s=16,alpha=0.75)
    for lbl,col,ctr in ref_types: ra.text(ctr[0],ctr[1]+0.38,lbl,ha="center",fontsize=7.5,color=col,fontweight="bold")
    ra.set_title("Reference Atlas\n(pre-annotated, e.g. 700k cells)",fontsize=9,fontweight="bold",color="#1A5276")
    ra.set_xticks([]); ra.set_yticks([]
                                     ); ra.spines[:].set_visible(False); ra.set_facecolor("#EBF5FB")

    # Query panel
    q_x=np.random.uniform(0.2,3,120); q_y=np.random.uniform(0.5,3.5,120)
    qa=fig.add_axes([0.37,0.1,0.24,0.56]); qa.scatter(q_x,q_y,c="#1A252F",s=16,alpha=0.45)
    qa.set_title("Your New Data\n(unannotated cells)",fontsize=9,fontweight="bold",color="#1A252F")
    qa.set_xticks([]); qa.set_yticks([]); qa.spines[:].set_visible(False); qa.set_facecolor("#F2F3F4")
    qa.text(1.5,0.1,"?",ha="center",fontsize=28,fontweight="bold",color="#5D6D7E",alpha=0.5)

    ax.text(6.35,2.7,"Integration\n+\nLabel transfer",ha="center",fontsize=10,fontweight="bold",color="#E67E22",
            bbox=dict(facecolor="#FDEBD0",edgecolor="#E67E22",boxstyle="round,pad=0.3",lw=1.5))
    ax.annotate("",xy=(7.1,2.7),xytext=(6.6,2.7),arrowprops=dict(arrowstyle="-|>",color="#E67E22",lw=2.5,mutation_scale=18))

    # Annotated output
    centres=[(t[2][0],t[2][1]) for t in ref_types]; ref_cols=[t[1] for t in ref_types]
    assigned=[ref_cols[np.argmin([((qx-c[0])**2+(qy-c[1])**2)**.5 for c in centres])] for qx,qy in zip(q_x,q_y)]
    oa=fig.add_axes([0.71,0.1,0.27,0.56]); oa.scatter(q_x,q_y,c=assigned,s=16,alpha=0.8)
    for lbl,col,ctr in ref_types:
        ms=[i for i,c in enumerate(assigned) if c==col]
        if ms: oa.text(np.mean([q_x[i] for i in ms]),np.mean([q_y[i] for i in ms])+0.3,lbl,
                       ha="center",fontsize=7.5,color=col,fontweight="bold")
    oa.set_title("Annotated Output\n(labels transferred)",fontsize=9,fontweight="bold",color="#27AE60")
    oa.set_xticks([]); oa.set_yticks([]); oa.spines[:].set_visible(False); oa.set_facecolor("#EAFAF1")
    return fig

# ══════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── Slide 1: Title ─────────────────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["dark"])
rect(sl, 0, Inches(2.2), SW, Inches(3.4), C["navy"])
txt(sl,"Cell-Type Annotation in Single-Cell Genomics",
    Inches(0.5),Inches(2.45),Inches(12.3),Inches(1.1),
    size=36,bold=True,color=C["white"],align=PP_ALIGN.CENTER)
txt(sl,"A gentle introduction for biology and chemistry audiences",
    Inches(0.5),Inches(3.5),Inches(12.3),Inches(0.6),
    size=18,color=C["lblue"],align=PP_ALIGN.CENTER,italic=True)

topics=["Single-Cell RNA-seq","Gene Expression Matrix","Marker Genes","Zero-Shot Annotation","Reference Mapping"]
for i,t in enumerate(topics):
    x=Inches(0.7+i*2.5)
    rect(sl,x,Inches(4.5),Inches(2.2),Inches(0.65),
         [C["blue"],C["orange"],C["purple"],C["green"],C["teal"]][i])
    txt(sl,t,x,Inches(4.52),Inches(2.2),Inches(0.62),
        size=11,bold=True,color=C["white"],align=PP_ALIGN.CENTER)

txt(sl,"Scripps Research Hackathon 2026",Inches(0.5),Inches(6.9),Inches(12.3),Inches(0.35),
    size=11,color=C["mid"],align=PP_ALIGN.CENTER)
footer(sl,1)

# ── Slide 2: Central Question ──────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"The Central Question",
       "A tissue contains many different cell types — but how do we tell them apart?",accent=C["navy"])
img(sl,fig_central_question(),Inches(0.3),Inches(0.9),Inches(12.7),Inches(6.1))
footer(sl,2)

# ── Slide 3: Gene Expression ───────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"What is Gene Expression?",
       "The molecular language that defines cell identity  |  Accessible for chemists and biologists",
       accent=C["purple"])
img(sl,fig_central_dogma(),Inches(0.3),Inches(0.9),Inches(12.7),Inches(6.1))
footer(sl,3)

# ── Slide 4: scRNA-seq Technology ─────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Single-Cell RNA Sequencing (scRNA-seq)",
       "Measuring gene activity in thousands of individual cells simultaneously",accent=C["teal"])
img(sl,fig_workflow(),Inches(0.3),Inches(0.9),Inches(12.7),Inches(6.0))
footer(sl,4)

# ── Slide 5: UMAP ─────────────────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"From Data Matrix to UMAP: Finding Cell Groups",
       "Dimensionality reduction collapses 20,000 genes → 2D map.  Similar cells cluster together.",
       accent=C["teal"])
img(sl,fig_umap_panels(),Inches(0.3),Inches(0.9),Inches(12.7),Inches(6.1))
footer(sl,5)

# ── Slide 6: Marker Genes ─────────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Marker Genes: The Molecular Fingerprints of Cell Identity",
       "Some genes are uniquely and highly expressed in specific cell types",accent=C["purple"])

img(sl,fig_dotplot(),Inches(0.2),Inches(0.92),Inches(7.8),Inches(5.2))

# Marker gene table (right)
trows=[["Gene(s)","Cell Type","Role"],
       ["CD3D / CD3E","T cells","T-cell receptor complex"],
       ["CD79A / MS4A1","B cells","B-cell receptor complex"],
       ["CD14 / LYZ","CD14+ Mono","Innate immunity / phagocytosis"],
       ["FCGR3A","FCGR3A+ Mono","FC receptor (non-classical)"],
       ["GNLY / NKG7","NK cells","Cytotoxic granule proteins"],
       ["FCER1A","Dendritic cells","IgE receptor"],
       ["PPBP / PF4","Megakaryocytes","Platelet factors"]]
ts=sl.shapes.add_table(len(trows),3,Inches(8.15),Inches(1.1),Inches(4.95),Inches(5.1))
tt=ts.table
for c,w in enumerate([1.45,1.45,2.05]): tt.columns[c].width=Inches(w)
row_bgs=["1B4F72","D6EAF8","FDEBD0","D6EAF8","FDEBD0","D6EAF8","FDEBD0","D6EAF8"]
for r,row in enumerate(trows):
    for c,val in enumerate(row):
        cell=tt.cell(r,c)
        set_cell_bg(cell,row_bgs[r])
        ctxt(cell,val,size=9 if r>0 else 10,bold=(r==0),
             color="FFFFFF" if r==0 else "1A252F")

footer(sl,6)

# ── Slide 7: Zero-Shot Annotation ─────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Method 1: Zero-Shot Annotation with Marker Genes",
       "Annotate without training data — use biological prior knowledge directly",accent=C["purple"])
img(sl,fig_zero_shot_flow(),Inches(0.3),Inches(0.88),Inches(12.7),Inches(6.1))
footer(sl,7)

# ── Slide 8: Reference Mapping ────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Method 2: Reference Mapping",
       "Compare your cells to a large, annotated atlas and borrow their labels",accent=C["teal"])
img(sl,fig_reference_map(),Inches(0.3),Inches(0.9),Inches(12.7),Inches(5.8))
callout(sl,"Tools used in this project: CellTypist (Claude) | Scanorama+KNN (Gemma) | Centroid dot-product (Llama)",
        Inches(0.3),Inches(6.85),Inches(12.7),Inches(0.4),C["dark"],C["lblue"])
footer(sl,8)

# ── Slide 9: Why It Matters ───────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Why Does Cell-Type Annotation Matter?",
       "Annotated cells unlock biological and clinical insights",accent=C["green"])

apps=[
    (0.25,1.05,C["blue"],  "Disease Research",       "Identify which cell types are expanded or\ndepleted in cancer, autoimmune disease, or infection"),
    (4.60,1.05,C["green"], "Drug Target Discovery",  "Find the exact cell type where a drug acts —\nor should act — in the body"),
    (8.95,1.05,C["orange"],"Immune Profiling",       "Track how T, B, NK cells change before and\nafter vaccination or treatment"),
    (0.25,3.55,C["purple"],"Cell Atlas Projects",    "Build maps of every cell type in the\nhuman body (Human Cell Atlas project)"),
    (4.60,3.55,C["red"],   "Biomarker Discovery",    "Identify rare cell populations that predict\npatient outcomes or disease progression"),
    (8.95,3.55,C["teal"],  "Drug Development",       "Understand off-target effects: which cell types\nare affected by a new compound"),
]
for l,t,color,title,desc in apps:
    rect(sl,Inches(l),Inches(t),Inches(4.1),Inches(2.2),color)
    txt(sl,title,Inches(l+0.08),Inches(t+0.1),Inches(3.94),Inches(0.55),
        size=13,bold=True,color=C["white"],align=PP_ALIGN.CENTER)
    txt(sl,desc,Inches(l+0.08),Inches(t+0.65),Inches(3.94),Inches(1.45),
        size=10,color=C["white"],align=PP_ALIGN.CENTER)

rect(sl,Inches(0.25),Inches(6.85),Inches(12.83),Inches(0.42),C["dark"])
txt(sl,"Without cell-type annotation, single-cell data is just 20,000 numbers per cell — annotation gives it biological meaning",
    Inches(0.35),Inches(6.87),Inches(12.63),Inches(0.38),
    size=11,bold=True,color=C["white"],align=PP_ALIGN.CENTER)
footer(sl,9)

# ── Slide 10: Summary Table ───────────────────────────────────────────────────
sl = blank(prs); bg(sl, C["light"])
header(sl,"Summary: Zero-Shot vs Reference Mapping",
       "Two complementary strategies — use both for the most reliable annotation",accent=C["dark"])

rows=[
    ["","Zero-Shot (Marker Genes)","Reference Mapping"],
    ["Core idea","Match cluster's top genes to known\nmarker gene lists","Align cells to a pre-annotated atlas\nand transfer labels"],
    ["Prior knowledge","Curated marker gene lists\n(literature or LLM knowledge)","A labelled reference dataset\n(e.g. Human Cell Atlas)"],
    ["Granularity","Cluster-level\n(all cells in a cluster get same label)","Cell-level\n(each cell annotated independently)"],
    ["Strengths","- No reference data required\n- Fast and interpretable\n- Any organism","- Cell-level resolution\n- Handles novel datasets\n- Large training sets"],
    ["Weaknesses","- Limited by cluster resolution\n- Misses rare cell types\n- Needs good markers","- Needs a good reference\n- Batch effects can mislead\n- Reference may be incomplete"],
    ["Accuracy (PBMC3k)","85.1% (Claude & Llama)\n76.4% (Gemma)","88.9% (Claude/CellTypist)\n77.0% (Gemma/Scanorama)\n73.2% (Llama/dot-product)"],
    ["Best when...","No reference available\nor working with novel tissue","High accuracy needed and\na well-matched reference exists"],
]
ts=sl.shapes.add_table(len(rows),3,Inches(0.25),Inches(0.95),Inches(12.83),Inches(6.15))
tt=ts.table
for c,w in enumerate([2.1,5.35,5.35]): tt.columns[c].width=Inches(w)
for r,row in enumerate(rows):
    for c,val in enumerate(row):
        cell=tt.cell(r,c)
        if r==0:
            bg_c=["1B4F72","2E86C1","1A5276"][c]
            set_cell_bg(cell,bg_c); ctxt(cell,val,size=12,bold=True,color="FFFFFF")
        else:
            bg_c=("EAECEE" if r%2==0 else "D5D8DC") if c==0 else \
                 (("EBF5FB" if r%2==0 else "D6EAF8") if c==1 else \
                  ("E8F8F5" if r%2==0 else "D5F5E3"))
            set_cell_bg(cell,bg_c); ctxt(cell,val,size=9.5,color="1A252F")

footer(sl,10)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "intro_cell_type_annotation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({Path(OUT).stat().st_size//1024} KB)")
